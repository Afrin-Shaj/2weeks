import streamlit as st
from pptx import Presentation
from PyPDF2 import PdfReader
from pdf2image import convert_from_bytes
import pytesseract
from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv
import google.generativeai as genai
import faiss
import numpy as np
import os
import uuid
import io
from PIL import Image
import pandas as pd

st.set_page_config(page_title="Startup RAG QA", layout="wide")


# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Initialize OCR if available
OCR_ENABLED = False
try:
    pytesseract.get_tesseract_version()
    OCR_ENABLED = True
except:
    st.warning("Tesseract OCR not installed. Image-based documents may not be processed fully.")

# Constants
UPLOAD_DIR = os.path.join(os.getcwd(), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Initialize models
@st.cache_resource
def load_models():
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    llm_model = genai.GenerativeModel('gemini-1.5-flash')
    return embedding_model, llm_model

embedding_model, llm_model = load_models()

# Helper functions
def extract_text_from_pdf(file_path):
    """Extract text from PDF, including image-based PDFs using OCR"""
    text_per_page = []
    with open(file_path, "rb") as f:
        pdf_bytes = f.read()
        
        # First try standard text extraction
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text = page.extract_text() or ""
            if len(text.strip()) > 30:  # If reasonable text found
                text_per_page.append(text)
            elif OCR_ENABLED:  # Try OCR for image-based pages
                try:
                    images = convert_from_bytes(pdf_bytes, first_page=page_num+1, last_page=page_num+1)
                    if images:
                        text = pytesseract.image_to_string(images[0])
                        text_per_page.append(text)
                    else:
                        text_per_page.append("")
                except Exception as e:
                    st.error(f"OCR failed for page {page_num+1}: {str(e)}")
                    text_per_page.append("")
            else:
                text_per_page.append("")
    return text_per_page

def extract_text_from_pptx(file_path):
    """Extract text from PPTX, including notes and slide text"""
    prs = Presentation(file_path)
    slides = []
    for slide in prs.slides:
        text = ""
        # Extract slide text
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        
        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                text += "Notes: " + notes_slide.notes_text_frame.text + "\n"
        
        slides.append(text.strip())
    return slides

def embed_chunks(chunks):
    return embedding_model.encode(chunks).astype("float32")

def create_faiss_index(chunks):
    vectors = embed_chunks(chunks)
    index = faiss.IndexFlatL2(vectors.shape[1])
    index.add(vectors)
    return index, vectors

def retrieve_similar_chunks(query, chunks, index, vectors, top_k=5, compare_mode=False):
    query_vec = embed_chunks([query])
    
    if compare_mode and len(st.session_state.file_map) > 1:
        # For comparison mode, retrieve balanced results from each document
        per_doc_top_k = max(1, top_k // len(st.session_state.file_map))
        top_chunks = []
        
        for doc_chunks in st.session_state.file_map.values():
            # Create temporary index for this document
            doc_vectors = embed_chunks(doc_chunks)
            doc_index = faiss.IndexFlatL2(doc_vectors.shape[1])
            doc_index.add(doc_vectors)
            
            # Retrieve from this document
            _, I = doc_index.search(query_vec, per_doc_top_k)
            top_chunks.extend([doc_chunks[i] for i in I[0]])
        return top_chunks[:top_k]
    else:
        # Standard retrieval
        _, I = index.search(query_vec, top_k)
        return [chunks[i] for i in I[0]]

def query_llm(context, question, compare_mode=False):
    """Query LLM with formatted context and instructions"""
    if compare_mode:
        prompt = f"""
# Role
You're a startup investment analyst comparing multiple pitch decks.

# Instructions
1. Compare the documents based on the question
2. Create a comparison table with these columns: Feature, {", ".join(st.session_state.file_map.keys())}
3. For each feature, extract values from each document
4. Highlight key differences using **bold**
5. Always cite the source document and slide number
6. If information is missing for a document, write "Not found"

# Context
{context}

# Question
{question}

# Response Format
- Start with a summary of key differences
- Then show the comparison table
- End with investment recommendations if applicable
"""
    else:
        prompt = f"""
# Role
You're a startup investment analyst reviewing pitch decks. 

# Instructions
1. Answer the question using ONLY the provided context
2. Always cite the source document and slide number
3. Format your answer using Markdown (headings, bullet points)
4. If information is missing, say "Information not found"

# Context
{context}

# Question
{question}

# Response Format
- Use **bold** for key metrics
- Use bullet points for lists
- Include document references in parentheses
"""
    response = llm_model.generate_content(prompt)
    return response.text

def extract_default_insights(doc_name, chunks):
    """Extract business model, revenue, and team insights for a specific document"""
    context = "\n\n".join([chunk.split("CONTENT:\n")[1] for chunk in chunks])
    
    business_model_prompt = f"Extract the business model from {doc_name}. How does the company make money?"
    revenue_prompt = f"Extract revenue model, pricing strategy, and financial projections from {doc_name}"
    team_prompt = f"Extract key team members, their expertise, and team structure from {doc_name}"
    
    with st.spinner(f"🔍 Analyzing {doc_name}..."):
        business_model = llm_model.generate_content(
            f"Based on this context:\n{context}\n\n{business_model_prompt}"
        ).text
    
        revenue = llm_model.generate_content(
            f"Based on this context:\n{context}\n\n{revenue_prompt}"
        ).text
    
        team = llm_model.generate_content(
            f"Based on this context:\n{context}\n\n{team_prompt}"
        ).text
    
    return {
        "Business Model": business_model,
        "Revenue Insights": revenue,
        "Team Insights": team
    }

# Streamlit UI
def main():
    # st.set_page_config(page_title="Startup RAG QA", layout="wide")
    st.title("📊 Startup Pitch Deck Analyzer")
    
    # Initialize session state
    if "all_chunks" not in st.session_state:
        st.session_state.all_chunks = []
        st.session_state.file_map = {}
        st.session_state.index = None
        st.session_state.vectors = None
        st.session_state.default_insights = {}
    
    # File uploader
    uploaded_files = st.file_uploader(
        "Upload pitch decks (PDF/PPTX)", 
        type=["pdf", "pptx"], 
        accept_multiple_files=True
    )
    
    # Process uploaded files
    if uploaded_files:
        with st.expander("File Processing Status", expanded=True):
            processed_files = set()
            if st.session_state.all_chunks:
                processed_files = {name for name in st.session_state.file_map.keys()}
            
            for file in uploaded_files:
                status = "✅ Processed" if file.name in processed_files else "🔄 Processing"
                st.write(f"{status}: {file.name}")
        
        # Process new files
        with st.spinner("Processing documents..."):
            for file in uploaded_files:
                if file.name not in st.session_state.file_map:
                    filename = f"{uuid.uuid4()}_{file.name}"
                    file_path = os.path.join(UPLOAD_DIR, filename)
                    with open(file_path, "wb") as f:
                        f.write(file.getbuffer())
                    
                    if file.name.endswith(".pdf"):
                        chunks = extract_text_from_pdf(file_path)
                    else:
                        chunks = extract_text_from_pptx(file_path)
                    
                    # Tag chunks with document and slide info
                    tagged_chunks = [
                        f"DOCUMENT: {file.name} | SLIDE: {i+1} | CONTENT:\n{chunk}" 
                        for i, chunk in enumerate(chunks)
                    ]
                    st.session_state.all_chunks.extend(tagged_chunks)
                    st.session_state.file_map[file.name] = tagged_chunks
            
            # Create/update FAISS index
            if st.session_state.all_chunks:
                st.session_state.index, st.session_state.vectors = create_faiss_index(
                    st.session_state.all_chunks
                )
    
    # Query section at the top
    if st.session_state.get("all_chunks"):
        st.subheader("🔍 Ask Questions")
        question = st.text_area("Enter your question", 
                              height=100,
                              placeholder="e.g., What's the revenue model? Compare business models across decks")
        
        col1, col2 = st.columns([1, 3])
        
        with col1:
            compare_mode = st.checkbox("Compare across documents", 
                                      disabled=len(st.session_state.file_map) < 2,
                                      help="Enable to compare features across multiple pitch decks")
            ask_btn = st.button("Analyze Documents")
        
        if ask_btn and question:
            # Retrieve relevant chunks
            top_chunks = retrieve_similar_chunks(
                question, 
                st.session_state.all_chunks,
                st.session_state.index,
                st.session_state.vectors,
                top_k=8,
                compare_mode=compare_mode
            )
            
            # Generate context
            context = "\n\n---\n\n".join(top_chunks)
            
            # Query LLM
            answer = query_llm(context, question, compare_mode=compare_mode)
            
            # Display results
            st.subheader("📈 Analysis Results")
            st.markdown(answer, unsafe_allow_html=True)
            
            with st.expander("📄 Reference Slides Used"):
                for i, chunk in enumerate(top_chunks):
                    doc_name = chunk.split(" | ")[0].replace("DOCUMENT: ", "")
                    slide_num = chunk.split(" | ")[1].replace("SLIDE: ", "")
                    content = " | ".join(chunk.split(" | ")[2:]).replace("CONTENT:\n", "")
                    
                    st.subheader(f"Reference #{i+1}: {doc_name} - Slide {slide_num}")
                    st.text(content)
    
    # Extract and display insights for each document
    if uploaded_files and st.session_state.file_map:
        st.markdown("---")
        st.subheader("🔑 Key Insights from Pitch Decks")
        
        # Extract insights if not already done
        if not st.session_state.default_insights:
            st.session_state.default_insights = {}
            for doc_name, chunks in st.session_state.file_map.items():
                st.session_state.default_insights[doc_name] = extract_default_insights(doc_name, chunks)
        
        # Display insights in tabs
        tabs = st.tabs([f"{doc_name}" for doc_name in st.session_state.file_map.keys()])
        
        for idx, (doc_name, insights) in enumerate(st.session_state.default_insights.items()):
            with tabs[idx]:
                st.markdown(f"### {doc_name}")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("#### Business Model")
                    st.markdown(insights["Business Model"])
                    
                    st.markdown("#### Revenue Insights")
                    st.markdown(insights["Revenue Insights"])
                
                with col2:
                    st.markdown("#### Team Insights")
                    st.markdown(insights["Team Insights"])
    
    # Sample questions
    st.markdown("---")
    st.subheader("💡 Sample Questions to Ask")
    
    if len(st.session_state.file_map) > 1:
        st.markdown("**Multi-Deck Comparisons**")
        st.markdown("- Compare the business models across decks")
        st.markdown("- Compare the target markets of these startups")
        st.markdown("- Compare the founding teams' experience")
        st.markdown("- Compare the revenue projections")
        st.markdown("- Compare the competitive advantages")
    else:
        st.markdown("**Single Deck Questions**")
        st.markdown("- What's the revenue model?")
        st.markdown("- Who is the target audience?")
        st.markdown("- What makes this product unique?")
        st.markdown("- What problem does this solve?")
        st.markdown("- What's the customer acquisition strategy?")

if __name__ == "__main__":
    main()